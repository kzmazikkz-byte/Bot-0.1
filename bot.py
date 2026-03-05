import telebot
from telebot import types
from pptx import Presentation
from pptx.util import Inches
import requests
import os
import random
import openai
from io import BytesIO
from PIL import Image
import time

# ----------------------
# Настройки токенов
# ----------------------
TOKEN = os.environ["TOKEN"]
OPENAI_KEY = os.environ["OPENAI_KEY"]
openai.api_key = OPENAI_KEY

bot = telebot.TeleBot(TOKEN)
user_data = {}

# ----------------------
# Функция прогресса
# ----------------------
def send_progress(message, steps=10, total_time=15, text_prefix="Генерирую презентацию…"):
    msg = bot.send_message(message.chat.id, f"{text_prefix} 0%")
    for i in range(1, steps+1):
        time.sleep(total_time/steps)
        percent = int(i/steps*100)
        try:
            bot.edit_message_text(f"{text_prefix} {percent}%", chat_id=message.chat.id, message_id=msg.message_id)
        except:
            pass
    return msg

# ----------------------
# Старт
# ----------------------
@bot.message_handler(commands=['start'])
def start(message):
    kb = types.ReplyKeyboardMarkup(resize_keyboard=True)
    kb.add("🦈 Создать презентацию")
    kb.add("✏ Редактировать слайды")
    bot.send_message(
        message.chat.id,
        "Привет! Я Shark AI v4. Создаю профессиональные презентации с помощью ИИ.",
        reply_markup=kb
    )

# ----------------------
# Создание презентации
# ----------------------
@bot.message_handler(func=lambda m: m.text=="🦈 Создать презентацию")
def create_presentation(message):
    msg = bot.send_message(message.chat.id, "📚 Напишите тему презентации:")
    bot.register_next_step_handler(msg, choose_style)

def choose_style(message):
    user_data[message.chat.id] = {"topic": message.text}
    kb = types.ReplyKeyboardMarkup(resize_keyboard=True)
    kb.add("🌙 Dark", "☀ Light", "💜 Neon")
    msg = bot.send_message(message.chat.id, "🎨 Выберите стиль презентации:", reply_markup=kb)
    bot.register_next_step_handler(msg, choose_slides)

def choose_slides(message):
    user_data[message.chat.id]["style"] = message.text
    msg = bot.send_message(message.chat.id, "📑 Сколько слайдов? (1-10)")
    bot.register_next_step_handler(msg, generate_presentation)

# ----------------------
# Генерация презентации
# ----------------------
def generate_presentation(message):
    try:
        count = int(message.text)
        if not 1 <= count <= 10:
            raise ValueError
    except:
        bot.send_message(message.chat.id, "❌ Введите число от 1 до 10.")
        return

    topic = user_data[message.chat.id]["topic"]
    style = user_data[message.chat.id]["style"]

    # Общий прогресс
    total_progress_msg = send_progress(message, steps=5, total_time=5, text_prefix="⚡ Подготавливаю презентацию…")

    # Прогресс генерации текста
    text_progress_msg = bot.send_message(message.chat.id, "🧠 Генерация текста слайдов через ИИ: 0%")
    for i in range(1, 11):
        time.sleep(0.5)
        try:
            bot.edit_message_text(f"🧠 Генерация текста слайдов через ИИ: {i*10}%", 
                                  chat_id=message.chat.id, 
                                  message_id=text_progress_msg.message_id)
        except:
            pass

    # Новый синтаксис openai>=1.0.0
    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Сделай презентацию на тему '{topic}' из {count} слайдов. Дай текст для каждого слайда в виде списка."}]
        )
        slides_text = response.choices[0].message["content"].split("\n")
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ Ошибка генерации текста: {e}")
        return

    # Создание презентации
    prs = Presentation()
    for i in range(count):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Заголовок
        try:
            title = slide.shapes.title
            title.text = f"{topic} – слайд {i+1}"
        except:
            pass

        # Текст
        try:
            body = slide.placeholders[1]
            body.text = slides_text[i] if i < len(slides_text) else topic
        except:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            textbox.text = slides_text[i] if i < len(slides_text) else topic

        # Картинка
        try:
            img_url = f"https://source.unsplash.com/800x600/?{topic}"
            img_data = requests.get(img_url).content
            img_stream = BytesIO(img_data)
            img = Image.open(img_stream)
            temp_file = f"/tmp/temp_{i}.png"
            img.save(temp_file)
            slide.shapes.add_picture(temp_file, Inches(6), Inches(2), height=Inches(3))
        except:
            pass

    # Сохраняем файл
    file_name = f"/tmp/shark_v4_{random.randint(1000,9999)}.pptx"
    prs.save(file_name)

    # Отправка презентации
    try:
        with open(file_name, "rb") as f:
            bot.send_document(message.chat.id, f)
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ Ошибка при отправке презентации: {e}")

    # Удаление сообщений прогресса
    try:
        bot.delete_message(message.chat.id, total_progress_msg.message_id)
        bot.delete_message(message.chat.id, text_progress_msg.message_id)
    except:
        pass

# ----------------------
# Запуск бота
# ----------------------
bot.infinity_polling()
